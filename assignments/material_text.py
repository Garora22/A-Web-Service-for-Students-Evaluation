"""Text extraction for different course material file types."""

from __future__ import annotations

from io import BytesIO
from typing import BinaryIO

from .pdf_text import extract_text_from_pdf


def _ext(name: str) -> str:
    name = (name or "").lower()
    if "." not in name:
        return ""
    return "." + name.rsplit(".", 1)[-1]


def extract_text_from_pptx(file_obj: BinaryIO) -> str:
    """
    Extract text from PPTX slides (text boxes only).
    Raises RuntimeError if python-pptx is not installed.
    """
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - env dependent
        raise RuntimeError(
            "PPTX support requires 'python-pptx'. Install with: pip install python-pptx"
        ) from exc

    raw = file_obj.read()
    if hasattr(file_obj, "seek"):
        file_obj.seek(0)

    prs = Presentation(BytesIO(raw))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                text = (shape.text or "").strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


def extract_text_from_plain_text(file_obj: BinaryIO) -> str:
    raw = file_obj.read()
    if hasattr(file_obj, "seek"):
        file_obj.seek(0)
    try:
        return raw.decode("utf-8", errors="ignore")
    except Exception:  # noqa: BLE001
        return ""


def extract_text_from_material(filename: str, file_obj: BinaryIO) -> tuple[str, str]:
    """
    Returns (text, kind) where kind is a short string like 'pdf'/'pptx'/'text'.
    Raises ValueError for unsupported types.
    """
    ext = _ext(filename)
    if ext == ".pdf":
        return extract_text_from_pdf(file_obj), "pdf"
    if ext == ".pptx":
        return extract_text_from_pptx(file_obj), "pptx"
    if ext in {".txt", ".md"}:
        return extract_text_from_plain_text(file_obj), "text"
    if ext == ".ppt":
        raise ValueError("Unsupported PowerPoint format .ppt. Please upload .pptx or export to PDF.")
    raise ValueError(f"Unsupported file type: {ext or '(no extension)'}")

